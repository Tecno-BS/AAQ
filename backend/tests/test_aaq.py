"""
Script para probar AAQ y su funcionamiento.

Las gráficas en PowerPoint (Insertar > Gráfico) no son imágenes: son objetos nativos.
Para que el LLM analice la diapositiva completa (título + texto + gráfico), se puede
exportar cada diapositiva como PNG con PowerPoint (solo Windows, requiere Office).

Uso: Levantar API -> desde backend: python tests/test_aaq.py
"""

import sys
import tempfile
from pathlib import Path

import requests
from docx import Document
from pptx import Presentation  # type: ignore[import]

BASE_URL = "http://localhost:8000"

def ask_path(prompt: str, default: str | None = None) -> Path:
    while True:
        raw = input(f"{prompt}" + (f" [{default}]" if default else "") + ": ").strip()
        if not raw and default:
            raw = default
        
        p = Path(raw)
        if p.exists():
            return p
        print(f"Ruta no válida o archivo no existe: {raw}")

def extract_context_from_docx(docx_path: Path) -> dict:
    #Extrae el texto del DOCX y construye el ResearchContext
    doc = Document(docx_path)
    paras = [p.text for p in doc.paragraphs if p.text]
    full_text = "\n".join(paras)

    profile = "Experto en medicamentos genéricos"
    background = full_text[:2500]
    business_question = (
        "¿Tiene La Santé la capacidad de expandirse hacia nuevas categorías y segmentos"
        "sin comprometer la coherencia de su marca?"
    )

    return {
        "profile" : profile,
        "background" : background,
        "business_question" : business_question,
        #Luego se pide el tipo de estudio válido al usuario
        "study_type" : None,
        "segments" : ["Hombres y mujeres 18+", "Estratos 2-6"],
        "sample" : "350 encuestas, margen 5.2%",
        "significance_threshold" : None,
        "models" : ["Modelo Aaker", "Driver Analysis"],
        "measurements" : [],
        "strategic_purposes" : ["Gestionar estrategia de salud de marca La Santé"],
        "qualitative_study" : None,
    }


def choose_study_type() -> str:
    #Le pide al backend la lista de tipos de estudio y deja que el usuario elija uno

    try:
        r = requests.get(f"{BASE_URL}/study-types", timeout=10)
        r.raise_for_status()
        data = r.json()
        types = data.get("study_types", [])
    except Exception as e:
        print(f"No se pudo obtener /study-types ({(e)}). Usando un valor fijo.")
        return "BET - Salud de marca"

    if not types:
        print("La lista de tipos de estudio está vacía. Usando un valor fijo.")
        return "BET - Salud de marca"

    print("\nTipos de estudio disponibles:")
    for idx, t in enumerate(types, start=1):
        print(f"  {idx}. {t}")
    
    while True:
        raw = input("Elige el numero del tipo de estudio: ").strip()

        if raw.isdigit():
            i = int(raw)
            if 1 <= i <= len(types):
                return types[i-1]
            print("Opción inválida. Intenta de nuevo.")

def inspect_pptx_slides(pptx_path: Path) -> list[int]:
    #Muestra las diapositivas y deja elegir cuáles analizar. 
    prs = Presentation(pptx_path)
    selected = []

    print("\nDiapositivas encontradas en el PPTX:")
    for i, slide in enumerate(prs.slides, start=1):
        #Obtiene el titulo
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text or ""
        #Contar imágenes
        img_count = 0
        for sh in slide.shapes:
            if hasattr(sh, "image"):
                img_count += 1
        print(f"  {i}: '{title[:60]}) (imagenes: {img_count})")
    
    raw = input("\nIndica que diapositivas quieres analizar (ej: 2,3,5) o ENTER para todas: ").strip()

    if not raw:
        return list(range(1, len(prs.slides) + 1))
    
    parts = [p.strip() for p in raw.split(",") if p.strip()]
    for p in parts:
        if p.isdigit():
            idx = int(p)
            if 1 <= idx <= len(prs.slides):
                selected.append(idx)
    
    if not selected:
        print("No seleccionaste ninguna valida se usaran todas.")
        return list(range(1, len(prs.slides) + 1))
    return selected

def export_slides_to_png_powerpoint(
    pptx_path: Path, selected_slides: list[int]
) -> list[tuple[bytes, str]] | None:
    """
    Exporta las diapositivas seleccionadas como PNG usando PowerPoint (solo Windows).

    Así se obtiene la diapositiva completa tal cual: título, texto y gráficos nativos
    (los hechos con Insertar > Gráfico), no solo imágenes incrustadas.

    Requiere: Microsoft PowerPoint instalado, pip install pywin32
    Retorna: list[(bytes, "png")] o None si falla (ej. sin Office / no Windows).
    """
    if sys.platform != "win32":
        print("export_slides_to_png_powerpoint: no es Windows (sys.platform != 'win32').")
        return None
    try:
        import win32com.client  # pywin32
    except ImportError:
        print(
            "export_slides_to_png_powerpoint: no se pudo importar win32com.client. "
            "Asegúrate de instalar pywin32 en este entorno (pip install pywin32)."
        )
        return None

    pptx_abs = str(pptx_path.resolve())
    results: list[tuple[bytes, str]] = []

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        try:
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            # En algunas versiones de Office no se permite ocultar la ventana.
            # Preferimos dejarla visible antes que fallar.
            try:
                ppt.Visible = True
            except Exception:
                pass
            pres = ppt.Presentations.Open(pptx_abs, WithWindow=False)
        except Exception as e:
            print(
                "export_slides_to_png_powerpoint: error al crear instancia de PowerPoint "
                f"o al abrir la presentación:\n{e!r}"
            )
            return None

        try:
            for slide_num in selected_slides:
                if slide_num < 1 or slide_num > pres.Slides.Count:
                    continue
                out_png = tmp / f"slide_{slide_num}.png"
                pres.Slides(slide_num).Export(str(out_png), "PNG")
                if out_png.exists():
                    results.append((out_png.read_bytes(), "png"))
            pres.Close()
        finally:
            ppt.Quit()

    return results if results else None


def extract_chart_data_from_pptx(
    pptx_path: Path, selected_slides: list[int]
) -> list[dict]:
    """
    Extrae datos estructurados de los gráficos nativos de PowerPoint
    (Insertar > Gráfico) en las diapositivas seleccionadas.

    Retorna una lista de dicts con:
    - slide_index
    - title
    - chart_type
    - series: [{name, points: [{category, value}]}]
    """
    prs = Presentation(pptx_path)
    charts: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide_idx not in selected_slides:
            continue

        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text or ""

        for shape in slide.shapes:
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                continue

            chart = shape.chart

            info: dict = {
                "slide_index": slide_idx,
                "title": title,
                "chart_type": str(getattr(chart, "chart_type", "")),
                "series": [],
            }

            categories: list[str | None] = []
            try:
                if chart.plots and chart.plots[0].categories is not None:
                    categories = [c.label for c in chart.plots[0].categories]
            except Exception:
                categories = []

            for series in chart.series:
                serie_data: dict = {
                    "name": getattr(series, "name", None),
                    "points": [],
                }
                for idx, point in enumerate(series.points):
                    try:
                        value = point.value
                    except Exception:
                        continue
                    category = categories[idx] if idx < len(categories) else None
                    serie_data["points"].append(
                        {
                            "category": category,
                            "value": value,
                        }
                    )

                info["series"].append(serie_data)

            if info["series"]:
                charts.append(info)

    return charts


def summarize_chart_data_for_context(charts: list[dict]) -> str:
    """
    Construye un resumen en texto plano de los datos de los gráficos,
    pensado para añadirse al background del ResearchContext.
    """
    if not charts:
        return ""

    lines: list[str] = []
    lines.append(
        "Información cuantitativa extraída automáticamente de los gráficos del PPTX seleccionado:"
    )

    for chart in charts:
        slide_idx = chart.get("slide_index")
        title = chart.get("title") or "Sin título"
        chart_type = chart.get("chart_type") or "Tipo de gráfico desconocido"
        lines.append(f"- Diapositiva {slide_idx}: '{title}' ({chart_type})")

        for serie in chart.get("series", []):
            serie_name = serie.get("name") or "Serie sin nombre"
            points_desc: list[str] = []
            for p in serie.get("points", []):
                cat = p.get("category") or "Sin categoría"
                val = p.get("value")
                points_desc.append(f"{cat}={val}")
            if points_desc:
                lines.append(f"  Serie '{serie_name}': " + "; ".join(points_desc))

    return "\n".join(lines)


def maybe_ask_qualitative_study() -> str | None:
    #Pregunta si hay estudio cuali anterior
    ans = input("\nEl estudio tiene un informe cualitativo anterior? (s/n): ").strip().lower()
    if ans not in ("s", "si", "sí"):
        return None
    
    pdf_path = ask_path("Ruta del PDF del estudio cualitativo")
    return f"Estudio cualitativo previo adjunto en: {pdf_path}"
    
def main():
    print("=== Prueba interactiva AAQ con archivos reales ===\n")

    #1. Contexto
    default_docx = str(Path("tests/3352/contexto agente IA.docx"))
    docx_path = ask_path("Ruta del DOCX de contexto", default=default_docx)
    context = extract_context_from_docx(docx_path)

    #1.1 Tipo de estudio fijo desde back
    study_type = choose_study_type()
    context["study_type"] = study_type
    print(f"\nTipo de estudio seleccionado: {study_type}")

    #1.2 Estudio cualitativo previo (opcional)
    qualitative_info = maybe_ask_qualitative_study()
    if qualitative_info:
        context["qualitative_study"] = qualitative_info

    #2. PPTX y seleccionar diapositivas 
    default_pptx = str(Path("tests/3352") / "I - 3359 -MAI2_Salud marca + Streaching Consumidor- V2.pptx")
    pptx_path = ask_path("\n Ruta del PPTX con las gráficas", default_pptx)

    selected_slides = inspect_pptx_slides(pptx_path)
    print(f"Diapositivas seleccionadas: {selected_slides}")

    #3. Extraer datos numéricos de los gráficos del PPTX y añadirlos al contexto
    print("\nExtrayendo datos numéricos de los gráficos del PPTX...")
    charts_data = extract_chart_data_from_pptx(pptx_path, selected_slides)
    if charts_data:
        summary = summarize_chart_data_for_context(charts_data)
        print(f"Se extrajeron datos de {len(charts_data)} gráfico(s).")
        # Enriquecemos el background con esta información cuantitativa
        context["background"] = (context.get("background") or "") + "\n\n" + summary
    else:
        print("No se encontraron gráficos con datos estructurados en las diapositivas seleccionadas.")

    #4. Obtener diapositivas como imágenes (diapositiva completa = título + texto + gráfico nativo)
    print("\nObteniendo diapositivas como imagen (diapositiva completa)...")
    imgs = export_slides_to_png_powerpoint(pptx_path, selected_slides) or []
    files = []

    if imgs:
        print(f"Exportadas {len(imgs)} diapositivas con PowerPoint (PNG).")
        for i, (blob, ext) in enumerate(imgs):
            ext = (ext or "png").lower()
            mime = f"image/{ext}" if ext in ("png", "jpg", "jpeg", "webp") else "application/octet-stream"
            files.append(("files", (f"slide_{i}.{ext}", blob, mime)))
    else:
        print(
            "PowerPoint no disponible o falló. No se generarán imágenes; "
            "el análisis usará solo el contexto y los datos numéricos extraídos."
        )

    #5. Crear estudio (ya con contexto enriquecido con datos de gráficos)
    print("\nCreando estudio via API /studies...")
    r = requests.post(f"{BASE_URL}/studies", json=context)
    if r.status_code != 201:
        print(f"Error al crear estudio: {r.status_code}\n{r.text}")
        return

    study = r.json()
    study_id = study["id"]
    print(f"Estudio creado con ID: {study_id}")

    # Subir gráficas solo si hay archivos válidos (imágenes)
    if files:
        print(f"Subiendo gráficas a /studies/{study_id}/charts ...")
        r = requests.post(f"{BASE_URL}/studies/{study_id}/charts", files=files)

        if r.status_code != 201:
            print(f"Error al subir graficas: {r.status_code}\n{r.text}")
            return
        
        charts = r.json().get("charts", [])
        print(f"Graficas registradas en BD: {len(charts)}")
    else:
        print("No se subieron gráficas porque no se generaron imágenes de diapositiva.")

    #6. Realizar análisis solo si hay al menos una gráfica asociada
    if files:
        print(f"Lanzando análisis en /studies/{study_id}/analyze ...")
        r = requests.post(f"{BASE_URL}/studies/{study_id}/analyze")
        if r.status_code != 202:
            print(f"Error al iniciar analisis: {r.status_code}\n{r.text}")
            return
        print("Analisis iniciado. El reporte estará disponible cuando finalice.")
    else:
        print(
            "No se puede iniciar el análisis porque no hay ninguna gráfica subida "
            "para este estudio. Necesitas al menos una diapositiva exportada a imagen."
        )

    print("\n=== Siguentes pasos ===")
    print("1) Espera 1-2 minutos mientras corre el pipeline")
    print(f'2) Consulta el reporte con: \n   curl "{BASE_URL}/reports?study_id={study_id}"')
    print(f"3) O ve a Swagger UI: {BASE_URL}/docs y prueba GET /reports con ese study_id.")

if __name__ == "__main__":
    main()